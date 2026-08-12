# ingest.py — Turning a file the CFO adds into something the agent can use.
# Two kinds of file arrive here: tabular (csv/tsv/xlsx/xls/parquet), which becomes
# a DataFrame, and documents (pdf/docx/pptx/txt/...), which become text.
#
# Raw extracted document text is poor input for the agent — PowerPoint comes out
# as shape fragments, Word tables collapse into loose lines, PDFs carry page
# furniture. So documents go through a second, deliberate step: a small
# converter agent (a cheap model, one non-streaming call) that normalises the
# extraction into clean markdown with the figures preserved verbatim. That runs
# once per file — results are cached by content hash — so the main agent always
# reads the tidy version and never pays for the conversion twice, and no answer
# ever stalls waiting on a conversion mid-turn.
#
# This module is a leaf by construction, like budget.py and cash.py: it imports
# pandas and the standard library, and nothing from this app at module level.
# The one app dependency — agent.get_client — is imported INSIDE convert_document
# for two reasons: it keeps ingest importable with no API key and no SDK on the
# path (which is what makes tests/test_ingest.py fixture-free), and it breaks a
# real cycle, since agent -> tools -> uploads -> ingest.
#
# Per-format readers are imported inside their extractor for the same reason
# export.py imports openpyxl inside workbook(): a checkout that predates the
# requirements bump loses one file format rather than failing to import the app.
#
# Everything here returns values or raises ValueError with a message meant to be
# shown to the user verbatim; nothing here talks to the UI.

from __future__ import annotations

import hashlib
import io
import json
import os
from pathlib import Path

import pandas as pd

DATA_DIR = Path(__file__).resolve().parent.parent / "data"
UPLOAD_DIR = DATA_DIR / "uploads"
CACHE_DIR = UPLOAD_DIR / "cache"

TABULAR_EXTS = {"csv", "tsv", "xlsx", "xls", "parquet"}
DOCUMENT_EXTS = {"pdf", "docx", "pptx", "txt", "md", "json", "log", "yaml", "yml", "xml"}
IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp"}
SUPPORTED_EXTS = TABULAR_EXTS | DOCUMENT_EXTS | IMAGE_EXTS

# The converter runs on a cheap model: it is a reformatting job, not analysis.
# Pinned to the same id as agent.REASONING_SUMMARY_MODEL so the app has one Haiku
# version rather than two that drift apart — tests/test_ingest.py asserts it.
CONVERTER_MODEL = "claude-haiku-4-5-20251001"

MAX_EXTRACT_CHARS = 200_000     # cap on text handed to the converter / the agent
MAX_PREVIEW_ROWS = 50           # rows of a chat-attached table inlined for the agent
MIN_PDF_TEXT_CHARS = 200        # below this a PDF is treated as scanned (see extract_text)

# Returned instead of text when a PDF yields almost nothing — it is image-only,
# so the caller should send the native base64 document block and let the model
# read the pages directly rather than pass on an empty extraction.
SCANNED_PDF = "__scanned_pdf__"

CONVERTER_SYSTEM_PROMPT = """You convert raw text extracted from a business document into clean \
markdown for a finance agent to read. You are a reformatter, not an analyst.

Rules:
- Preserve every number, date, currency amount, percentage and label EXACTLY as it appears. Never \
round, recompute, reformat or infer a figure. If a value is garbled in the extraction, keep it as-is.
- Reconstruct tables as markdown tables, written left-flush at the start of the line with a header \
row and a |---| separator. This is the main thing you add: extracted tables usually arrive as loose \
lines.
- Keep the document's headings and section order. Use markdown headings.
- Emit the blocks in exactly the order they appear in the input. Never move a paragraph, \
table or heading relative to the others — the input is already in document order, and a \
sentence written after a table must stay after it.
- Drop page furniture: page numbers, repeated headers/footers, navigation text, slide-number \
artefacts.
- For slides, keep one section per slide titled with the slide title, and include speaker notes \
under a "Notes:" line when present.
- Add nothing of your own: no summary, no commentary, no analysis, no preamble. Output only the \
converted document.
"""


def extension(filename: str) -> str:
    return (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()


def classify(filename: str) -> str:
    """"tabular" | "document" | "image" | "unsupported" — from the extension.

    This binary split matters everywhere downstream: datasets are QUERIED (filter,
    group, aggregate) and documents are READ (whole markdown body, one tool call).
    The Data ingestion page renders them as two sections for the same reason —
    mixing them would misrepresent what the agent can do with each.
    """
    ext = extension(filename)
    if ext in TABULAR_EXTS:
        return "tabular"
    if ext in DOCUMENT_EXTS:
        return "document"
    if ext in IMAGE_EXTS:
        return "image"
    return "unsupported"


def supported_note() -> str:
    """Human-readable list of accepted extensions, for error messages and the UI."""
    return ", ".join("." + e for e in sorted(SUPPORTED_EXTS))


# --------------------------------------------------------------------------
# Tabular files
# --------------------------------------------------------------------------

def read_table(filename: str, raw: bytes) -> pd.DataFrame:
    """Parse a tabular upload into a DataFrame. Raises ValueError with a
    user-facing message on anything unreadable."""
    ext = extension(filename)
    buf = io.BytesIO(raw)
    try:
        if ext == "parquet":
            df = pd.read_parquet(buf)
        elif ext in ("xlsx", "xls"):
            # Only the first sheet becomes the dataset; sheet_names is reported
            # by the caller so the user knows the rest was left behind.
            df = pd.read_excel(buf, sheet_name=0)
        elif ext == "tsv":
            df = pd.read_csv(buf, sep="\t")
        else:
            df = pd.read_csv(buf)
    except Exception as exc:
        raise ValueError(f"Could not read '{filename}' as a table: {exc}") from exc
    if df.empty or not len(df.columns):
        raise ValueError(f"'{filename}' contains no rows or no columns.")
    # Normalise column names so the query tools can address them predictably.
    df.columns = [str(c).strip() for c in df.columns]
    return _stringify_dates(df)


def _stringify_dates(df: pd.DataFrame) -> pd.DataFrame:
    """Store date columns as strings, matching the curated datasets.

    Excel and Parquet hand back real datetimes, which json-serialise to epoch
    milliseconds — meaningless to the model. Month-grained columns become
    'YYYY-MM' (the convention every period helper in tools.py, budget.py and
    cash.py matches on), anything finer becomes 'YYYY-MM-DD'.
    """
    for col in df.columns:
        if not pd.api.types.is_datetime64_any_dtype(df[col]):
            continue
        values = df[col]
        month_grained = bool(((values.dt.day == 1) & (values.dt.normalize() == values)).all())
        df[col] = values.dt.strftime("%Y-%m" if month_grained else "%Y-%m-%d")
    return df


def sheet_names(filename: str, raw: bytes) -> list[str]:
    """Sheet names of an Excel upload ([] for anything else) — so the UI can say
    which sheet became the dataset."""
    if extension(filename) not in ("xlsx", "xls"):
        return []
    try:
        return list(pd.ExcelFile(io.BytesIO(raw)).sheet_names)
    except Exception:
        return []


def summarise_table(df: pd.DataFrame, filename: str) -> str:
    """Inline form of a chat-attached table: shape, dtypes and the first rows as
    CSV. Truncation is STATED so the model never assumes it saw everything —
    the same discipline as query_budget_data's `truncated` flag."""
    head = df.head(MAX_PREVIEW_ROWS)
    dtypes = ", ".join(f"{c} ({df[c].dtype})" for c in df.columns)
    lines = [f"Attached table: {filename}",
             f"{len(df)} rows x {len(df.columns)} columns.",
             f"Columns: {dtypes}", ""]
    if len(df) > MAX_PREVIEW_ROWS:
        lines.append(f"First {MAX_PREVIEW_ROWS} of {len(df)} rows "
                     f"(the rest was not included — ask the user to add this file "
                     f"on the Data ingestion page to query all of it):")
    else:
        lines.append("Full contents:")
    lines.append("```csv")
    lines.append(head.to_csv(index=False).strip())
    lines.append("```")
    return "\n".join(lines)


# --------------------------------------------------------------------------
# Document text extraction
# --------------------------------------------------------------------------

def _extract_pdf(raw: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(raw))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            pages.append(f"[page {i}]\n{text.strip()}")
    return "\n\n".join(pages)


def _docx_table_text(table) -> str:
    """A Word table, pipe-separated so the converter can rebuild it as markdown."""
    rows = []
    for row in table.rows:
        cells = [c.text.strip().replace("\n", " ") for c in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "[table]\n" + "\n".join(rows) if rows else ""


def _extract_docx(raw: bytes) -> str:
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = docx.Document(io.BytesIO(raw))
    # Walk the body in document order rather than document.paragraphs followed by
    # document.tables: that shortcut loses the interleaving, so a sentence written
    # after a table ends up before it and the converter has no way to recover it.
    parts = []
    for child in document.element.body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            text = Paragraph(child, document).text.strip()
            if text:
                parts.append(text)
        elif tag == "tbl":
            text = _docx_table_text(Table(child, document))
            if text:
                parts.append(text)
    return "\n\n".join(parts)


def _extract_pptx(raw: bytes) -> str:
    from pptx import Presentation
    presentation = Presentation(io.BytesIO(raw))
    slides = []
    for i, slide in enumerate(presentation.slides, start=1):
        parts = [f"[slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(f"Notes: {notes}")
        if len(parts) > 1:
            slides.append("\n".join(parts))
    return "\n\n".join(slides)


def _extract_plain(raw: bytes) -> str:
    return raw.decode("utf-8", errors="replace").strip()


def extract_text(filename: str, raw: bytes) -> str:
    """Deterministic first pass: document bytes -> raw text.

    Returns the SCANNED_PDF sentinel for an image-only PDF so the caller can fall
    back to the native document block — "we could not read text out of this" and
    "this document is empty" are different states, and pretending the first is the
    second ships an empty extraction the model has no way to question.

    Raises ValueError if the file can't be opened at all.
    """
    ext = extension(filename)
    try:
        if ext == "pdf":
            text = _extract_pdf(raw)
            if len(text.strip()) < MIN_PDF_TEXT_CHARS:
                return SCANNED_PDF
        elif ext == "docx":
            text = _extract_docx(raw)
        elif ext == "pptx":
            text = _extract_pptx(raw)
        else:
            text = _extract_plain(raw)
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError(f"Could not read '{filename}': {exc}") from exc
    if not text.strip():
        raise ValueError(f"No readable text found in '{filename}'.")
    if len(text) > MAX_EXTRACT_CHARS:
        # Never silently dropped: the marker is what tells the model it is reading
        # part of a file.
        text = text[:MAX_EXTRACT_CHARS] + "\n\n[truncated — file is longer than " \
                                          f"{MAX_EXTRACT_CHARS} characters]"
    return text


# --------------------------------------------------------------------------
# The converter agent
# --------------------------------------------------------------------------

def convert_document(filename: str, raw_text: str, model: str = CONVERTER_MODEL) -> str:
    """Normalise raw extracted text into clean markdown via a small model call.

    Falls back to the raw text when there is no API key, or the call fails for any
    reason — the same deterministic-fallback stance as alerts.py's narrative and
    budgetplan.py's templated_narrative: a degraded document beats a failed upload.
    """
    from .agent import get_client  # local: no key needed to import, and breaks the cycle
    if not os.environ.get("ANTHROPIC_API_KEY"):
        return raw_text
    try:
        response = get_client().messages.create(
            model=model,
            max_tokens=16000,
            system=[{"type": "text", "text": CONVERTER_SYSTEM_PROMPT,
                     "cache_control": {"type": "ephemeral"}}],
            messages=[{"role": "user", "content":
                       f"Convert this extracted text from '{filename}' to clean markdown.\n\n"
                       f"{raw_text}"}],
        )
        converted = "".join(b.text for b in response.content if b.type == "text").strip()
        return converted or raw_text
    except Exception:
        return raw_text


def _cache_path(raw: bytes) -> Path:
    return CACHE_DIR / f"{hashlib.sha256(raw).hexdigest()}.md"


def convert_cached(filename: str, raw: bytes, model: str = CONVERTER_MODEL) -> str:
    """extract_text + convert_document, cached by content hash.

    Keyed on the BYTES, not the name, so re-attaching the same file — or uploading
    it under a different name — is free, and the conversion happens exactly once
    per distinct file content, ever. Returns the SCANNED_PDF sentinel straight
    through.
    """
    text = extract_text(filename, raw)
    if text == SCANNED_PDF:
        return SCANNED_PDF
    path = _cache_path(raw)
    if path.exists():
        try:
            return path.read_text(encoding="utf-8")
        except OSError:
            pass
    converted = convert_document(filename, text, model=model)
    header = f"Document: {filename}\n\n"
    if not converted.startswith("Document:"):
        converted = header + converted
    try:
        CACHE_DIR.mkdir(parents=True, exist_ok=True)
        tmp = path.with_suffix(".md.tmp")
        tmp.write_text(converted, encoding="utf-8")
        tmp.replace(path)
    except OSError:
        pass  # cache is an optimisation, never a hard requirement
    return converted


def cache_key(raw: bytes) -> str:
    """Content hash used as the cache file name — stored on upload records so a
    document's converted markdown can be found again."""
    return hashlib.sha256(raw).hexdigest()


def cached_markdown(key: str) -> str | None:
    """Markdown for a stored document, or None if it is missing.

    The key is an OPAQUE file name. For a converted upload it is the content hash
    of the source bytes; for markdown the user owns (a saved answer, or a document
    they have edited) uploads.py uses the record's id instead — so this directory
    holds both cached conversions and user-owned text, and neither this function
    nor read_document nor _remove_files needs to know which it is holding.
    """
    path = CACHE_DIR / f"{key}.md"
    if not path.exists():
        return None
    try:
        return path.read_text(encoding="utf-8")
    except OSError:
        return None


def describe_table(df: pd.DataFrame, filename: str, sheets: list[str] | None = None) -> str:
    """One-line description for a dataset registered from an upload — mirrors the
    style of the built-in descriptions in tools.DATASETS."""
    bits = [f"Uploaded from {filename}", f"{len(df)} rows", f"{len(df.columns)} columns"]
    if "month" in df.columns:
        months = pd.to_datetime(df["month"], errors="coerce").dropna()
        if len(months):
            bits.append(f"covering {months.min():%Y-%m} to {months.max():%Y-%m}")
    if sheets and len(sheets) > 1:
        bits.append(f"sheet '{sheets[0]}' of {len(sheets)}")
    return ", ".join(bits) + "."


def json_safe(df: pd.DataFrame, rows: int) -> list[dict]:
    """First N rows as NaN-safe JSON records, for the dataset preview endpoint —
    the same job tools._records does for tool results."""
    return json.loads(df.head(rows).to_json(orient="records"))
